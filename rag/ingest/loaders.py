"""
Optimized Document Loaders for RAG Pipeline
============================================
Key Improvements:
1. Better table detection with multi-level fallbacks
2. Smarter text normalization (preserves structure when needed)
3. Robust error handling with detailed logging
4. Memory-efficient processing
5. Better integration with semantic chunker
"""

from typing import List, Optional, Any, Dict, Tuple
import re
import string
from pathlib import Path
from dataclasses import dataclass, field
from rag.logging import logger

fitz: Optional[Any] = None
HAS_PYMUPDF = False

try:
    import fitz
    HAS_PYMUPDF = True
    logger.info("✓ PyMuPDF loaded")
except ImportError:
    logger.warning("✗ PyMuPDF not available - PDF extraction limited")


try:
    from rag.ocr.pdf_ocr import (
        extract_text_from_page,
        is_available as ocr_is_available,
        get_status as ocr_get_status
    )
    HAS_OCR = True
    if ocr_is_available():
        logger.info("✓ OCR module loaded and available")
    else:
        logger.warning("⚠ OCR module loaded but engine not available (missing dependencies or models)")
        HAS_OCR = False
except ImportError as e:
    HAS_OCR = False
    logger.warning(f"✗ OCR module not available: {e}")
except Exception as e:
    HAS_OCR = False
    logger.warning(f"✗ OCR module error: {e}", exc_info=True)

    def extract_text_from_page(page: Any, **kwargs) -> Tuple[str, Optional[Dict]]:
        return ("", None)

    def ocr_is_available() -> bool:
        return False

    def ocr_get_status() -> dict:
        return {"available": False}

# Language detection
try:
    from langdetect import detect, DetectorFactory
    DetectorFactory.seed = 0  # For consistent results
    HAS_LANGDETECT = True
except ImportError:
    HAS_LANGDETECT = False
    logger.warning("✗ langdetect not available - language detection disabled")


# Core Document Model
@dataclass
class Document:
    """Document with text content and metadata"""
    text: str
    metadata: Dict[str, Any] = field(default_factory=dict)

    def __post_init__(self):
        if not isinstance(self.text, str):
            raise TypeError("Document text must be a string")
        if not isinstance(self.metadata, dict):
            raise TypeError("Document metadata must be a dictionary")

    def __len__(self) -> int:
        return len(self.text)

    def __repr__(self) -> str:
        preview = self.text[:50] + "..." if len(self.text) > 50 else self.text
        return f"Document(text='{preview}', metadata={self.metadata})"


def _make_document(text: str, **metadata) -> Document:
    """Factory function to create documents"""
    return Document(text=text, metadata=metadata)


# Text Quality & Normalization
_PRINTABLE = set(string.printable)
_WHITESPACE_PATTERN = re.compile(r'[ \t]+')
_MULTIPLE_NEWLINES = re.compile(r'\n{3,}')
_MULTIPLE_SPACES = re.compile(r' {2,}')


def _normalize_text(text: str, preserve_structure: bool = False) -> str:
    """
    Smart text normalization.
    
    Args:
        text: Raw text to normalize
        preserve_structure: If True, keep newlines and table formatting
    
    Returns:
        Cleaned text
    """
    if not text:
        return ""
    
    if preserve_structure:
        text = _WHITESPACE_PATTERN.sub(' ', text)
        text = _MULTIPLE_NEWLINES.sub('\n\n', text)
        text = _MULTIPLE_SPACES.sub(' ', text)
    else:
        text = re.sub(r'\s+', ' ', text)
    
    return text.strip()


def _printable_ratio(text: str) -> float:
    """Calculate ratio of printable/valid characters"""
    if not text:
        return 0.0
    valid_chars = sum(1 for c in text if c in _PRINTABLE or ord(c) > 127)
    return valid_chars / len(text)


def _detect_language(text: str) -> Optional[str]:
    """
    Detect language from text sample.
    
    Args:
        text: Text to detect language from
    
    Returns:
        Language code (e.g., 'vi', 'en', 'ja') or None if detection fails
    """
    if not HAS_LANGDETECT or not text or len(text.strip()) < 50:
        return None
    
    try:
        sample = text[:1000] if len(text) > 1000 else text
        lang = detect(sample)
        
        lang_map = {
            'vi': 'vi',
            'en': 'en',
            'ja': 'japan',
            'ko': 'korean',
            'zh-cn': 'ch',
            'zh-tw': 'chinese_cht',
            'th': 'th',
        }
        
        return lang_map.get(lang, lang)
    except Exception as e:
        logger.debug(f"Language detection failed: {e}")
        return None


def _extract_table_metadata(text: str) -> Dict[str, Any]:
    """
    Extract table metadata from text.
    
    Args:
        text: Text that may contain tables
    
    Returns:
        Dictionary with table metadata
    """
    metadata = {
        "has_table": False,
        "table_count": 0,
        "table_rows": 0,
        "table_columns": 0
    }
    
    if not text or " | " not in text:
        return metadata
    
    lines = text.split('\n')
    table_rows = 0
    max_columns = 0
    in_table = False
    
    for line in lines:
        if "[TABLE]" in line:
            in_table = True
            metadata["table_count"] += 1
            continue
        if "[/TABLE]" in line:
            in_table = False
            continue
        
        if in_table or (" | " in line and len(line.split(" | ")) >= 2):
            table_rows += 1
            columns = len(line.split(" | "))
            max_columns = max(max_columns, columns)
    
    if table_rows > 0:
        metadata["has_table"] = True
        metadata["table_rows"] = table_rows
        metadata["table_columns"] = max_columns
    
    return metadata


def _tag_table_content(text: str) -> str:
    """
    Add semantic markers for table content.
    
    Args:
        text: Text that may contain tables (with | separators)
    
    Returns:
        Text with [TABLE] markers around table sections
    """
    if not text or " | " not in text:
        return text
    
    lines = text.split('\n')
    result = []
    in_table = False
    table_lines = []
    empty_line_count = 0
    max_empty_lines_in_table = 1
    
    for line in lines:
        is_table_row = " | " in line and len(line.split(" | ")) >= 2
        
        if is_table_row:
            if not in_table:
                result.append("[TABLE]")
                in_table = True
                empty_line_count = 0
            table_lines.append(line)
        else:
            if in_table:
                if not line.strip():
                    empty_line_count += 1
                    if empty_line_count <= max_empty_lines_in_table:
                        table_lines.append(line)
                        continue
                    else:
                        result.extend(table_lines)
                        result.append("[/TABLE]")
                        table_lines = []
                        in_table = False
                        empty_line_count = 0
                        result.append(line)
                else:
                    result.extend(table_lines)
                    result.append("[/TABLE]")
                    table_lines = []
                    in_table = False
                    empty_line_count = 0
                    result.append(line)
            else:
                result.append(line)
    
    if in_table:
        result.extend(table_lines)
        result.append("[/TABLE]")
    
    return "\n".join(result)


def _is_usable_text(text: str, min_len: int = 100, min_ratio: float = 0.6) -> bool:
    """
    Check if extracted text is usable for RAG.
    
    Args:
        text: Text to validate
        min_len: Minimum character length
        min_ratio: Minimum printable character ratio
    
    Returns:
        True if text meets quality standards
    """
    text = text.strip()
    if len(text) < min_len:
        return False
    if _printable_ratio(text) < min_ratio:
        return False
    return True


# =====================================================
# PDF Layout Extraction
# =====================================================

def _merge_pymupdf_and_ocr(
    pymupdf_text: Optional[str],
    ocr_text: Optional[str],
    use_smart_merge: bool = True
) -> Optional[str]:
    """
    Intelligently merge results from PyMuPDF and OCR.
    
    
    Args:
        pymupdf_text: Text from PyMuPDF (may be None)
        ocr_text: Text from OCR (may be None)
        use_smart_merge: Use line-by-line alignment and printable ratio
    
    Returns:
        Merged text or best available text
    """
    if not pymupdf_text and not ocr_text:
        return None
    
    if not pymupdf_text:
        return ocr_text.strip() if ocr_text else None
    if not ocr_text:
        return pymupdf_text.strip()
    
    if not use_smart_merge:
        # Simple merge
        pymupdf_len = len(pymupdf_text.strip())
        ocr_len = len(ocr_text.strip())
        
        if pymupdf_len > ocr_len * 1.5:
            return pymupdf_text.strip() if ocr_len < 50 else f"{pymupdf_text}\n\n{ocr_text}".strip()
        if ocr_len > pymupdf_len * 1.5:
            return ocr_text.strip() if pymupdf_len < 50 else f"{ocr_text}\n\n{pymupdf_text}".strip()
        return f"{pymupdf_text}\n\n{ocr_text}".strip()
    
    # Smart merge: line-by-line with printable ratio check
    pymupdf_lines = pymupdf_text.strip().split('\n')
    ocr_lines = ocr_text.strip().split('\n')
    
    merged_lines = []
    pymupdf_idx = 0
    ocr_idx = 0
    
    while pymupdf_idx < len(pymupdf_lines) or ocr_idx < len(ocr_lines):
        if pymupdf_idx >= len(pymupdf_lines):
            merged_lines.extend(ocr_lines[ocr_idx:])
            break
        if ocr_idx >= len(ocr_lines):
            merged_lines.extend(pymupdf_lines[pymupdf_idx:])
            break
        
        pymupdf_line = pymupdf_lines[pymupdf_idx].strip()
        ocr_line = ocr_lines[ocr_idx].strip()
        
        if not pymupdf_line:
            pymupdf_idx += 1
            continue
        if not ocr_line:
            ocr_idx += 1
            continue
        
        # Check printable ratio - prefer OCR for low quality PyMuPDF lines
        pymupdf_ratio = _printable_ratio(pymupdf_line)
        
        if pymupdf_ratio < 0.5 and len(ocr_line) > len(pymupdf_line) * 0.7:
            # PyMuPDF line has low quality, prefer OCR
            merged_lines.append(ocr_line)
            ocr_idx += 1
            pymupdf_idx += 1  # Skip corresponding PyMuPDF line
        elif len(pymupdf_line) > len(ocr_line) * 1.3 and pymupdf_ratio > 0.7:
            # PyMuPDF line is longer and high quality
            merged_lines.append(pymupdf_line)
            pymupdf_idx += 1
            if len(ocr_line) < 20:  # Skip short OCR line
                ocr_idx += 1
        else:
            # Similar quality - prefer longer or combine
            if abs(len(pymupdf_line) - len(ocr_line)) < 10:
                # Very similar, use PyMuPDF (usually more accurate)
                merged_lines.append(pymupdf_line)
            else:
                # Use longer one
                merged_lines.append(pymupdf_line if len(pymupdf_line) > len(ocr_line) else ocr_line)
            pymupdf_idx += 1
            ocr_idx += 1
    
    return "\n".join(merged_lines).strip()


def _extract_pdf_with_layout(page) -> str:
    """
    Extract PDF text preserving layout structure.
    
    Strategy:
    1. Get text blocks with position info
    2. Group blocks into lines based on Y-position
    3. Detect table rows vs prose
    4. Format with appropriate separators
    5. Fallback to dict method if extraction fails
    
    Args:
        page: PyMuPDF page object
    
    Returns:
        Extracted text with preserved structure
    """
    try:
        blocks = page.get_text("blocks")
        if not blocks:
            return ""
        
        text_blocks = []
        for block in blocks:
            if len(block) < 5:
                continue
            
            text = block[4].strip()
            if not text:
                continue
            
            x0, y0, x1, y1 = block[0], block[1], block[2], block[3]
            
            text_blocks.append({
                "text": text,
                "x0": x0,
                "y0": y0,
                "x1": x1,
                "y1": y1,
                "center_x": (x0 + x1) / 2,
                "center_y": (y0 + y1) / 2,
                "width": x1 - x0,
                "height": y1 - y0
            })
        
        if not text_blocks:
            return ""
        
        text_blocks.sort(key=lambda b: (b["y0"], b["x0"]))
        
        lines = _group_blocks_into_lines(text_blocks)
        result = "\n".join(lines)
        
        if len(result.strip()) < 50:
            result = _extract_using_dict_method(page)
        
        # Add table tags for better semantic understanding
        result = _tag_table_content(result)
        
        return result.strip()
    
    except Exception as e:
        logger.error(f"Layout extraction failed: {e}")
        return ""


def _group_blocks_into_lines(
    text_blocks: List[Dict],
    line_tolerance: float = 5.0
) -> List[str]:
    """
    Group text blocks into lines based on vertical position.
    
    Args:
        text_blocks: List of block dictionaries with position data
        line_tolerance: Max Y-distance to consider blocks on same line
    
    Returns:
        List of formatted line strings
    """
    lines = []
    current_line = []
    current_y = None
    
    for block in text_blocks:
        y_pos = block["center_y"]
        
        if current_y is None or abs(y_pos - current_y) <= line_tolerance:
            current_line.append(block)
            if current_y is None:
                current_y = y_pos
            else:
                current_y = (current_y + y_pos) / 2
        else:
            if current_line:
                line_text = _format_line_with_table_structure(current_line)
                if line_text:
                    lines.append(line_text)
            
            current_line = [block]
            current_y = y_pos
    
    if current_line:
        line_text = _format_line_with_table_structure(current_line)
        if line_text:
            lines.append(line_text)
    
    return lines


def _format_line_with_table_structure(blocks: List[Dict]) -> str:
    """
    Format a line with intelligent table detection.
    
    Detection heuristics (structure-based only):
    - Consistent gaps between columns
    - Horizontal alignment
    - Presence of numbers/currency (as supporting indicators)
    - Multiple distinct columns
    
    Note: Does not rely on keywords as they are too specific and unreliable.
    
    Args:
        blocks: Text blocks in reading order
    
    Returns:
        Formatted line with | separators for tables, spaces for prose
    """
    if not blocks:
        return ""
    
    if len(blocks) == 1:
        return blocks[0]["text"]
    
    blocks.sort(key=lambda b: b["x0"])
    
    gaps = []
    for i in range(len(blocks) - 1):
        gap = blocks[i + 1]["x0"] - blocks[i]["x1"]
        if gap > 0:
            gaps.append(gap)
    
    is_table_row = _detect_table_row(blocks, gaps)
    
    if is_table_row:
        return " | ".join(b["text"] for b in blocks)
    else:
        parts = []
        for i, block in enumerate(blocks):
            if i > 0:
                gap = block["x0"] - blocks[i - 1]["x1"]
                separator = " | " if gap > 20 else " "
                parts.append(separator)
            parts.append(block["text"])
        return "".join(parts)


def _detect_table_row(blocks: List[Dict], gaps: List[float]) -> bool:
    """
    Multi-heuristic table row detection.
    
    Args:
        blocks: Text blocks in the line
        gaps: Horizontal gaps between blocks
    
    Returns:
        True if line appears to be a table row
    """
    if len(blocks) < 2:
        return False
    
    if gaps:
        avg_gap = sum(gaps) / len(gaps)
        
        if len(gaps) > 1:
            variance = sum((g - avg_gap) ** 2 for g in gaps) / len(gaps)
            std_dev = variance ** 0.5
            coefficient_of_variation = std_dev / max(avg_gap, 1)
            
            if avg_gap > 10 and coefficient_of_variation < 0.6:
                return True
        else:
            if avg_gap > 20:
                return True
    
    if len(blocks) >= 2:
        y_positions = [b["center_y"] for b in blocks]
        y_avg = sum(y_positions) / len(y_positions)
        y_variance = sum((y - y_avg) ** 2 for y in y_positions) / len(y_positions)
        
        if y_variance < 15:
            x_positions = [b["x0"] for b in blocks]
            x_spread = max(x_positions) - min(x_positions)
            if x_spread > 80:
                return True
    
    # Additional indicators for table detection (structure-based only)
    if len(blocks) >= 3:
        has_numbers = any(re.search(r'\d+', b["text"]) for b in blocks)
        has_currency = any(
            re.search(r'[$€£¥₹₫₩]|USD|EUR|GBP|VND|JPY', b["text"])
            for b in blocks
        )
        
        # Only use structure-based indicators (numbers/currency are common in tables)
        # Don't rely on keywords as they are too specific and unreliable
        if (has_numbers or has_currency) and len(blocks) >= 3:
            return True
    
    return False


def _extract_using_dict_method(page) -> str:
    """
    Fallback extraction using PyMuPDF's dict method.
    More granular than blocks, useful when block extraction fails.
    
    Args:
        page: PyMuPDF page object
    
    Returns:
        Extracted text
    """
    try:
        dict_text = page.get_text("dict")
        if not dict_text or "blocks" not in dict_text:
            return ""
        
        lines = []
        for block in dict_text["blocks"]:
            if "lines" not in block:
                continue
            
            line_parts = []
            for line in block["lines"]:
                if "spans" not in line:
                    continue
                
                span_texts = [
                    span.get("text", "").strip()
                    for span in line["spans"]
                ]
                span_texts = [s for s in span_texts if s]
                
                if span_texts:
                    line_parts.append(" ".join(span_texts))
            
            if line_parts:
                lines.append(" | ".join(line_parts))
        
        return "\n".join(lines)
    
    except Exception as e:
        logger.error(f"Dict method extraction failed: {e}")
        return ""


# =====================================================
# PDF Loader
# =====================================================

def load_pdf(filepath: str) -> List[Document]:
    """
    Load PDF with PyMuPDF + PaddleOCR for optimal performance.
    
    Args:
        filepath: Path to PDF file
    
    Returns:
        List of Document objects (one per usable page)
    
    Raises:
        FileNotFoundError: If file doesn't exist
        ValueError: If no text can be extracted
    """
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"PDF not found: {filepath}")
    
    if not HAS_PYMUPDF:
        raise ValueError("PyMuPDF is required for PDF extraction")
    
    documents: List[Document] = []
    
    # Open PDF with PyMuPDF
    try:
        doc = fitz.open(filepath)
        total_pages = len(doc)
    except Exception as e:
        raise ValueError(f"Failed to open PDF: {e}")
    
    if total_pages == 0:
        raise ValueError(f"PDF has no pages: {filepath}")
    
    logger.info(f"Processing PDF: {path.name} ({total_pages} pages)")

    try:
        for page_num in range(1, total_pages + 1):
            pymupdf_text = None
            ocr_text = None
            extractor_name = None
            
            page = doc[page_num - 1]
            ocr_metadata = None
            
            try:
                pymupdf_text = _extract_pdf_with_layout(page)
                if not pymupdf_text or len(pymupdf_text.strip()) < 50:
                    raw = page.get_text("text")
                    pymupdf_text = _normalize_text(raw, preserve_structure=False)
                else:
                    pymupdf_text = _normalize_text(pymupdf_text, preserve_structure=True)
                
                if pymupdf_text and _is_usable_text(pymupdf_text, min_len=50):
                    extractor_name = "pymupdf"
            except Exception as e:
                logger.debug(f"PyMuPDF extraction failed for page {page_num}: {e}")
            
            if (not extractor_name or (pymupdf_text and len(pymupdf_text.strip()) < 100)) and HAS_OCR and ocr_is_available():
                try:
                    lang_hint = None
                    if pymupdf_text and len(pymupdf_text.strip()) > 50:
                        lang_hint = _detect_language(pymupdf_text)
                    
                    ocr_result = extract_text_from_page(
                        page,
                        dpi=300,  # High quality rendering
                        text_score=0.5,  # Filter low confidence
                        box_thresh=0.5,
                        lang_hint=lang_hint,
                        return_confidence=True
                    )
                    ocr_text, ocr_metadata = ocr_result
                    
                    if ocr_text and ocr_text.strip():
                        ocr_text = _normalize_text(ocr_text, preserve_structure=True)
                        ocr_text = _tag_table_content(ocr_text)
                        
                        if _is_usable_text(ocr_text, min_len=50):
                            if pymupdf_text and _is_usable_text(pymupdf_text, min_len=30):
                                merged_text = _merge_pymupdf_and_ocr(
                                    pymupdf_text,
                                    ocr_text,
                                    use_smart_merge=True
                                )
                                if merged_text and _is_usable_text(merged_text, min_len=50):
                                    extractor_name = "pymupdf+ocr"
                                    pymupdf_text = merged_text
                                else:
                                    extractor_name = "ocr"
                                    pymupdf_text = ocr_text
                            else:
                                extractor_name = "ocr"
                                pymupdf_text = ocr_text
                except Exception as e:
                    logger.debug(f"OCR failed for page {page_num}: {e}")
            
            if pymupdf_text and _is_usable_text(pymupdf_text, min_len=50):
                table_metadata = _extract_table_metadata(pymupdf_text)
                
                doc_metadata = {
                    "source": filepath,
                    "filename": path.name,
                    "page": page_num,
                    "total_pages": total_pages,
                    "doc_type": "pdf",
                    "extractor": extractor_name or "pymupdf",
                    "is_ocr": (extractor_name and "ocr" in extractor_name),
                    **table_metadata
                }
                
                if extractor_name and "ocr" in extractor_name and ocr_metadata:
                    doc_metadata["ocr_metadata"] = ocr_metadata
                
                documents.append(
                    _make_document(pymupdf_text, **doc_metadata)
                )
            else:
                logger.debug(f"Page {page_num} could not be extracted with any method")
    
    finally:
        doc.close()
    
    if documents:
        extractor_counts = {}
        for doc in documents:
            ext = doc.metadata.get("extractor", "unknown")
            extractor_counts[ext] = extractor_counts.get(ext, 0) + 1
        
        logger.info(
            f"✓ Extracted {len(documents)}/{total_pages} pages. "
            f"Methods: {extractor_counts}"
        )
        return documents

    raise ValueError(
        f"Cannot extract text from: {filepath}\n"
        f"PyMuPDF: {HAS_PYMUPDF}, OCR: {HAS_OCR and ocr_is_available()}"
    )


# =====================================================
# DOCX Loader
# =====================================================

def load_docx(filepath: str) -> List[Document]:
    """Load DOCX document"""
    try:
        from docx import Document as Docx
    except ImportError:
        raise ImportError("python-docx required: pip install python-docx")
    
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"DOCX not found: {filepath}")

    doc = Docx(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(paragraphs)
    text = _normalize_text(text, preserve_structure=False)

    if not _is_usable_text(text):
        logger.warning(f"No usable text in {path.name}")
        return []

    logger.info(f"✓ Extracted {len(paragraphs)} paragraphs from {path.name}")
    return [_make_document(
        text,
        source=filepath,
        filename=path.name,
        doc_type="docx",
        paragraph_count=len(paragraphs)
    )]


# =====================================================
# XLSX Loader
# =====================================================

def load_xlsx(filepath: str) -> List[Document]:
    """Load XLSX spreadsheet"""
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl required: pip install openpyxl")
    
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"XLSX not found: {filepath}")

    documents: List[Document] = []
    
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
    except Exception as e:
        logger.error(f"Failed to load workbook: {e}")
        return []

    for sheet_name in wb.sheetnames:
        try:
            ws = wb[sheet_name]
            rows = []

            for row in ws.iter_rows(values_only=True):
                cell_values = [str(c) for c in row if c is not None]
                if cell_values:
                    rows.append(" | ".join(cell_values))

            text = "\n".join(rows)
            text = _normalize_text(text, preserve_structure=True)
            
            if _is_usable_text(text):
                documents.append(
                    _make_document(
                        text,
                        source=filepath,
                        filename=path.name,
                        sheet=sheet_name,
                        doc_type="xlsx",
                        row_count=len(rows)
                    )
                )
                logger.info(f"✓ Sheet '{sheet_name}': {len(rows)} rows")
            else:
                logger.debug(f"Sheet '{sheet_name}' insufficient data")
                
        except Exception as e:
            logger.error(f"Failed to process sheet '{sheet_name}': {e}")
            continue

    wb.close()
    
    if not documents:
        logger.warning(f"No usable data in {path.name}")
    
    return documents


# =====================================================
# PPTX Loader
# =====================================================

def load_pptx(filepath: str) -> List[Document]:
    """Load PPTX presentation"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx required: pip install python-pptx")
    
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {filepath}")

    documents: List[Document] = []
    
    try:
        pres = Presentation(filepath)
    except Exception as e:
        logger.error(f"Failed to load presentation: {e}")
        return []

    for idx, slide in enumerate(pres.slides):
        texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                value = shape.text.strip()
                if value:
                    texts.append(value)

        text = "\n".join(texts)
        text = _normalize_text(text, preserve_structure=False)
        
        if _is_usable_text(text, min_len=50):
            documents.append(
                _make_document(
                    text,
                    source=filepath,
                    filename=path.name,
                    slide=idx + 1,
                    total_slides=len(pres.slides),
                    doc_type="pptx",
                    shape_count=len(texts)
                )
            )
            logger.debug(f"✓ Slide {idx + 1}: {len(texts)} shapes")
        else:
            logger.debug(f"Slide {idx + 1} insufficient text")

    if not documents:
        logger.warning(f"No usable text in {path.name}")
    
    logger.info(f"✓ Extracted {len(documents)}/{len(pres.slides)} slides")
    return documents


# =====================================================
# TXT Loader
# =====================================================

def load_txt(filepath: str) -> List[Document]:
    """Load plain text file with encoding detection"""
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"TXT not found: {filepath}")
    
    text = ""
    encodings = ["utf-8", "utf-16", "latin-1", "cp1252"]
    
    for enc in encodings:
        try:
            with open(filepath, encoding=enc) as f:
                text = f.read()
            logger.info(f"✓ Decoded {path.name} as {enc}")
            break
        except (UnicodeDecodeError, UnicodeError):
            continue
        except Exception as e:
            logger.error(f"Failed to read with {enc}: {e}")
            continue
    
    if not text:
        raise ValueError(f"Cannot decode file with any encoding: {filepath}")
    
    text = _normalize_text(text, preserve_structure=False)

    if not _is_usable_text(text):
        logger.warning(f"No usable text in {path.name}")
        return []

    logger.info(f"✓ Loaded {len(text)} chars from {path.name}")
    return [_make_document(
        text,
        source=filepath,
        filename=path.name,
        doc_type="txt",
        char_count=len(text)
    )]


# =====================================================
# Main Router
# =====================================================

SUPPORTED_EXTENSIONS = {
    "pdf": load_pdf,
    "docx": load_docx,
    "xlsx": load_xlsx,
    "pptx": load_pptx,
    "txt": load_txt,
}


def load(filepath: str) -> List[Document]:
    """
    Load document with automatic format detection.
    
    Supported: PDF, DOCX, XLSX, PPTX, TXT
    
    Args:
        filepath: Path to document
    
    Returns:
        List of Document objects
    
    Raises:
        FileNotFoundError: If file doesn't exist
        ValueError: If format unsupported or extraction fails
    
    Example:
        >>> docs = load("invoice.pdf")
        >>> for doc in docs:
        ...     print(f"Page {doc.metadata['page']}: {len(doc.text)} chars")
    """
    path = Path(filepath)
    
    if not path.exists():
        raise FileNotFoundError(f"File not found: {filepath}")
    
    ext = path.suffix.lower().lstrip('.')
    
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported format: .{ext}\n"
            f"Supported: {', '.join(f'.{e}' for e in SUPPORTED_EXTENSIONS.keys())}"
        )
    
    loader_func = SUPPORTED_EXTENSIONS[ext]
    
    try:
        return loader_func(str(path))
    except Exception as e:
        logger.error(f"Failed to load {path.name}: {e}")
        raise


def get_loader_info() -> Dict[str, Any]:
    """Get loader capabilities and status"""
    return {
        "supported_formats": list(SUPPORTED_EXTENSIONS.keys()),
        "dependencies": {
            "pymupdf": HAS_PYMUPDF,
            "ocr": HAS_OCR and ocr_is_available(),
        },
        "ocr_status": ocr_get_status() if HAS_OCR else {"available": False},
    }
